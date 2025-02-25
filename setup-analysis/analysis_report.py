from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re
from pathlib import Path
import textwrap
from typing import Optional

SSD_PATH = Path("/Volumes/Realtek 1Tb")


def export_to_ppt():
    """_summary_"""

    # Create presentation and get slide dimensions.
    prs = Presentation()
    width, height = prs.slide_width, prs.slide_height

    # Use desired slide layout.
    slide_layouts = {layout.name: layout for layout in prs.slide_layouts}
    layout_name = ""
    slide_layout = slide_layouts.get(layout_name, slide_layouts["Blank"])

    # Define margins and spacing.
    left_margin = Inches(0.2)
    top_margin = Inches(0.2)
    margins = (left_margin, top_margin)

    h_spacing = Inches(0.1)
    v_spacing = Inches(0.1)
    spacing = (h_spacing, v_spacing)

    # Calculate available space and determine image sizes for a grid of 3 columns x 2 rows.
    available_width = width - (2 * left_margin) - (2 * h_spacing)
    available_height = height - (2 * top_margin) - (1 * v_spacing)
    pic_width = available_width / 3
    pic_height = available_height / 2
    pic_dims = (pic_width, pic_height)

    def add_slide_nb(slide, slide_number):
        # Define dimensions for the slide number textbox.
        num_width = Inches(1)
        num_height = Inches(0.4)
        left_num = width - num_width - Inches(0.2)  # margin from right
        top_num = height - num_height - Inches(0.2)  # margin from bottom
        txBox = slide.shapes.add_textbox(left_num, top_num, num_width, num_height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_number)
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.RIGHT

    def add_images(images, slide_layout, margins, spacing, pic_dims, rows=2, cols=3):
        img_per_slide = rows * cols
        left_margin, top_margin = margins
        h_spacing, v_spacing = spacing
        pic_width, pic_height = pic_dims
        slide = None

        for i, image_path in enumerate(images):
            # * Every img_per_slide images, add a new slide
            if i % img_per_slide == 0:
                slide = prs.slides.add_slide(slide_layout)

            # * Calculate grid position.
            col = i % cols  # columns: 0, 1, 2
            row = (i % img_per_slide) // cols
            left = left_margin + col * (pic_width + h_spacing)
            top = top_margin + row * (pic_height + v_spacing)
            slide.shapes.add_picture(
                str(image_path), left, top, width=pic_width, height=pic_height
            )

    def add_section_slide(
        prs, text, font, size, alignment, wrap_width: Optional[int] = None
    ):
        slide_layout = prs.slide_layouts[6]  # using a known blank layout index
        slide = prs.slides.add_slide(slide_layout)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Create a textbox that spans the entire slide.
        txBox = slide.shapes.add_textbox(0, 0, slide_width, slide_height)
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        if wrap_width is not None:
            # Wrap the text to a given character width.
            text = "\n".join(textwrap.wrap(text, width=wrap_width))

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = size
        p.font.name = font
        p.alignment = alignment

    analyzed_data_dir = SSD_PATH / "PhD Data/experiment1-analysis"

    section_slide_params = {
        "font": "Aptos Display (Headings)",
        "size": Pt(54),
        "alignment": PP_ALIGN.CENTER,
    }
    # * HUMAN DATA - BEHAVIORAL
    section_name = "Human Data\nBehavioral"

    img_dir = analyzed_data_dir / "Lab/analyzed/group_lvl"
    add_section_slide(prs, section_name, **section_slide_params)

    images = sorted([f for f in img_dir.glob("*.png") if not f.name.startswith(".")])
    _images = [f for f in images if re.search("Accuracy", f.stem)]
    add_images(_images, slide_layout, margins, spacing, pic_dims, rows=1, cols=2)

    _images = [f for f in images if re.search("RT", f.stem)]
    add_images(_images, slide_layout, margins, spacing, pic_dims, rows=2, cols=2)

    # * Representational Similarity Analysis
    section_name = "Representational Similarity \nAnalysis"
    add_section_slide(
        prs, section_name, "Aptos Display (Headings)", Pt(60), PP_ALIGN.CENTER
    )

    # * HUMAN DATA - RSA
    img_dir = analyzed_data_dir / "Lab/analyzed/RSA-FRP-occipital-metric_correlation"
    images = sorted([f for f in img_dir.glob("*.png") if not f.name.startswith(".")])

    section_name = "Human Data\nRSA - Occipital - Item level"
    add_section_slide(prs, section_name, **section_slide_params)

    _images = [f for f in images if re.search(".*subj.*item_lvl", f.stem)]
    add_images(_images, slide_layout, margins, spacing, pic_dims)

    section_name = "Human Data\nRSA - Occipital - Pattern level"
    add_section_slide(prs, section_name, **section_slide_params)

    _images = [f for f in images if re.search(".*subj.*pattern_lvl", f.stem)]
    add_images(_images, slide_layout, margins, spacing, pic_dims)

    _images = [f for f in images if re.search(".*group", f.stem)]
    add_images(_images, slide_layout, margins, spacing, pic_dims)

    # * ANN DATA - RSA
    img_dir = analyzed_data_dir / "ANNs/analyzed/RSA-seq_tokens-metric_correlation"
    images = sorted([f for f in img_dir.glob("*.png") if not f.name.startswith(".")])

    section_name = "ANN Data\nRSA - Seq Tokens - Item level"
    add_section_slide(prs, section_name, **section_slide_params)

    _images = [f for f in images if re.search(".*item_lvl", f.stem)]
    add_images(_images, slide_layout, margins, spacing, pic_dims)

    section_name = "ANN Data\nRSA - Seq Tokens - Pattern level"
    add_section_slide(prs, section_name, **section_slide_params)

    _images = [f for f in images if re.search(".*pattern_lvl", f.stem)]
    add_images(_images, slide_layout, margins, spacing, pic_dims)

    # Save the presentation.
    prs.save("test2.pptx")
